import os, sys
import pandas as pd
from pptx import Presentation
import xlwings as xw
import tkinter as tk
from tkinter import filedialog, messagebox

#---------------------------------ITERATE TRHOUGH PPTX FILE---------------------------------#

def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for placeholder, replacement in replacements.items():
                if placeholder in run.text:
                    run.text = run.text.replace(placeholder, replacement)

def replace_text_in_table(table, replacements):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    for placeholder, replacement in replacements.items():
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, replacement)

#---------------------------------PROCESS OIC-AD DATA---------------------------------#

def process_data_ad():
    try:
        excel_path = AD_excel_entry.get()
        pptx_path = AD_pptx_entry.get()
        
        if not excel_path or not pptx_path:
            messagebox.showerror("Error", "Please select both Excel and PowerPoint files")
            return

        # Read data from Excel
        df = xw.Book(excel_path)
        
        # Read data from pptx
        prs = Presentation(pptx_path)
        
        # Read data from sheets
        dfass = df.sheets['adminloginsuccess table']
        dfnass = df.sheets['nonadminloginsuccess table']
        dfanys = df.sheets['anonymous login user_table']
        dfanyfs = df.sheets['anonymous loginfail src']
        dfcas = df.sheets['createuser user_table']
        dfeas = df.sheets['enableuser user_table']
        dfrps = df.sheets['resetpassword user_table']
        dfdas = df.sheets['disableuser user_table']

        # Iterate values in sheets
        dfas = dfass.range("A2:D11").value
        dfnas = dfnass.range("A2:D11").value
        dfany = dfanys.range("A2:C11").value
        dfanyf = dfanyfs.range("A2:C11").value
        dfca = dfcas.range("A2:C29").value
        dfea = dfeas.range("A2:C29").value
        dfrp = dfrps.range("A2:C11").value
        dfda = dfdas.range("A2:C11").value

        # Define your replacements (can be dynamically generated)
        replacements = {
    # Admin Login Failure
    # Non-admin Login Failure
    # Admin Login Success
        # Username:
        'asu1': str(dfas[0][0]), 'asu2': str(dfas[1][0]), 'asu3': str(dfas[2][0]), 'asu4': str(dfas[3][0]), 'asu5': str(dfas[4][0]), 'asu6': str(dfas[5][0]), 'asu7': str(dfas[6][0]), 'asu8': str(dfas[7][0]), 'asu9': str(dfas[8][0]), 'asu0': str(dfas[9][0]),
        # Source Address No:
        'asip1': str(dfas[0][1]), 'asip2': str(dfas[1][1]), 'asip3': str(dfas[2][1]), 'asip4': str(dfas[3][1]), 'asip5': str(dfas[4][1]), 'asip6': str(dfas[5][1]), 'asip7': str(dfas[6][1]), 'asip8': str(dfas[7][1]), 'asip9': str(dfas[8][1]), 'asip0': str(dfas[9][1]),
        # Source Address No:
        # Destination Host Name:
        'ashn1': str(dfas[0][1]), 'ashn2': str(dfas[1][1]), 'ashn3': str(dfas[2][1]), 'ashn4': str(dfas[3][1]), 'ashn5': str(dfas[4][1]), 'ashn6': str(dfas[5][1]), 'ashn7': str(dfas[6][1]), 'ashn8': str(dfas[7][1]), 'ashn9': str(dfas[8][1]), 'ashn0': str(dfas[9][1]),
        # count:
        'cas1': str(dfas[0][3])[:-2], 'cas2': str(dfas[1][3])[:-2], 'cas3': str(dfas[2][3])[:-2], 'cas4': str(dfas[3][3])[:-2], 'cas5': str(dfas[4][3])[:-2], 'cas6': str(dfas[5][3])[:-2], 'cas7': str(dfas[6][3])[:-2], 'cas8': str(dfas[7][3])[:-2], 'cas9': str(dfas[8][3])[:-2], 'cas0': str(dfas[9][3])[:-2],
    # Non-admin Login Success
        # Username:
        'ansu1': str(dfnas[0][0]), 'ansu2': str(dfnas[1][0]), 'ansu3': str(dfnas[2][0]), 'ansu4': str(dfnas[3][0]), 'ansu5': str(dfnas[4][0]), 'ansu6': str(dfnas[5][0]), 'ansu7': str(dfnas[6][0]), 'ansu8': str(dfnas[7][0]), 'ansu9': str(dfnas[8][0]), 'ansu0': str(dfnas[9][0]), 
        # Source Address:
        'ansip1': str(dfnas[0][1]), 'ansip2': str(dfnas[1][1]), 'ansip3': str(dfnas[2][1]), 'ansip4': str(dfnas[3][1]), 'ansip5': str(dfnas[4][1]), 'ansip6': str(dfnas[5][1]), 'ansip7': str(dfnas[6][1]), 'ansip8': str(dfnas[7][1]), 'ansip9': str(dfnas[8][1]), 'ansip0': str(dfnas[9][1]), 
        # Source Address Name:
        # Destination Host Name:
        'anshn1': str(dfnas[0][2]), 'anshn2': str(dfnas[1][2]), 'anshn3': str(dfnas[2][2]), 'anshn4': str(dfnas[3][2]), 'anshn5': str(dfnas[4][2]), 'anshn6': str(dfnas[5][2]), 'anshn7': str(dfnas[6][2]), 'anshn8': str(dfnas[7][2]), 'anshn9': str(dfnas[8][2]), 'anshn0': str(dfnas[9][2]), 
        # Count:
        'cans1': str(dfnas[0][3])[:-2], 'cans2': str(dfnas[1][3])[:-2], 'cans3': str(dfnas[2][3])[:-2], 'cans4': str(dfnas[3][3])[:-2], 'cans5': str(dfnas[4][3])[:-2], 'cans6': str(dfnas[5][3])[:-2], 'cans7': str(dfnas[6][3])[:-2], 'cans8': str(dfnas[7][3])[:-2], 'cans9': str(dfnas[8][3])[:-2], 'cans0': str(dfnas[9][3])[:-2], 
    # Anonymous Login Success
        # Source Address:
        'anys1': str(dfany[0][0]), 'anys2': str(dfany[1][0]), 'anys3': str(dfany[2][0]), 'anys4': str(dfany[3][0]), 'anys5': str(dfany[4][0]), 'anys6': str(dfany[5][0]), 'anys7': str(dfany[6][0]), 'anys8': str(dfany[7][0]), 'anys9': str(dfany[8][0]), 'anys0': str(dfany[9][0]), 
        # Source Address Name:
        # Destination Host Name:
        'anyhn1': str(dfnas[0][1]), 'anyhn2': str(dfnas[1][1]), 'anyhn3': str(dfnas[2][1]), 'anyhn4': str(dfnas[3][1]), 'anyhn5': str(dfnas[4][1]), 'anyhn6': str(dfnas[5][1]), 'anyhn7': str(dfnas[6][1]), 'anyhn8': str(dfnas[7][1]), 'anyhn9': str(dfnas[8][1]), 'anyhn0': str(dfnas[9][1]), 
        # Count:
        'cany1': str(dfany[0][2]), 'cany2': str(dfany[1][2]), 'cany3': str(dfany[2][2]), 'cany4': str(dfany[3][2]), 'cany5': str(dfany[4][2]), 'cany6': str(dfany[5][2]), 'cany7': str(dfany[6][2]), 'cany8': str(dfany[7][2]), 'cany9': str(dfany[8][2]), 'cany0': str(dfany[9][2]), 
    # Anonymous Login Failure
        # Source Address:
        'anyf1': str(dfanyf[0][0]), 'anyf2': str(dfanyf[1][0]), 'anyf3': str(dfanyf[2][0]), 'anyf4': str(dfanyf[3][0]), 'anyf5': str(dfanyf[4][0]), 'anyf6': str(dfanyf[5][0]), 'anyf7': str(dfanyf[6][0]), 'anyf8': str(dfanyf[7][0]), 'anyf9': str(dfanyf[8][0]), 'anyf0': str(dfanyf[9][0]), 
        # Source Address Name:
        # Destination Host Name:
        'anyhnf1': str(dfanyf[0][1]), 'anyhnf2': str(dfanyf[1][1]), 'anyhnf3': str(dfanyf[2][1]), 'anyhnf4': str(dfanyf[3][1]), 'anyhnf5': str(dfanyf[4][1]), 'anyhnf6': str(dfanyf[5][1]), 'anyhnf7': str(dfanyf[6][1]), 'anyhnf8': str(dfanyf[7][1]), 'anyhnf9': str(dfanyf[8][1]), 'anyhnf0': str(dfanyf[9][1]), 
        # Event Count:
        'canyf1': str(dfanyf[0][2]), 'canyf2': str(dfanyf[1][2]), 'canyf3': str(dfanyf[2][2]), 'canyf4': str(dfanyf[3][2]), 'canyf5': str(dfanyf[4][2]), 'canyf6': str(dfanyf[5][2]), 'canyf7': str(dfanyf[6][2]), 'canyf8': str(dfanyf[7][2]), 'canyf9': str(dfanyf[8][2]), 'canyf0': str(dfanyf[9][2]), 
    # Create Accunt
        # Username:
        'cacu1': str(dfca[0][0]), 'cacu2': str(dfca[1][0]), 'cacu3': str(dfca[2][0]), 'cacu4': str(dfca[3][0]), 'cacu5': str(dfca[4][0]), 'cacu6': str(dfca[5][0]), 'cacu7': str(dfca[6][0]), 'cacu8': str(dfca[7][0]), 'cacu9': str(dfca[8][0]), 'cacu0': str(dfca[9][0]),
        'cacua1': str(dfca[10][0]), 'cacua2': str(dfca[11][0]), 'cacua3': str(dfca[12][0]), 'cacua4': str(dfca[13][0]), 'cacua5': str(dfca[14][0]), 'cacua6': str(dfca[15][0]), 'cacua7': str(dfca[16][0]), 'cacua8': str(dfca[17][0]), 'cacua9': str(dfca[18][0]), 'cacua0': str(dfca[19][0]),
        'cacub1': str(dfca[20][0]), 'cacub2': str(dfca[21][0]), 'cacub3': str(dfca[22][0]), 'cacub4': str(dfca[23][0]), 'cacub5': str(dfca[24][0]), 'cacub6': str(dfca[25][0]), 'cacub7': str(dfca[26][0]), 'cacub8': str(dfca[27][0]),
        # Destination Username:
        'dunca1': str(dfca[0][1]), 'dunca2': str(dfca[1][1]), 'dunca3': str(dfca[2][1]), 'dunca4': str(dfca[3][1]), 'dunca5': str(dfca[4][1]), 'dunca6': str(dfca[5][1]), 'dunca7': str(dfca[6][1]), 'dunca8': str(dfca[7][1]), 'dunca9': str(dfca[8][1]), 'dunca0': str(dfca[9][1]),
        'duncaa1': str(dfca[10][1]), 'duncaa2': str(dfca[11][1]), 'duncaa3': str(dfca[12][1]), 'duncaa4': str(dfca[13][1]), 'duncaa5': str(dfca[14][1]), 'duncaa6': str(dfca[15][1]), 'duncaa7': str(dfca[16][1]), 'duncaa8': str(dfca[17][1]), 'duncaa9': str(dfca[18][1]), 'duncaa0': str(dfca[19][1]),
        'duncab1': str(dfca[20][1]), 'duncab2': str(dfca[21][1]), 'duncab3': str(dfca[22][1]), 'duncab4': str(dfca[23][1]), 'duncab5': str(dfca[24][1]), 'duncab6': str(dfca[25][1]), 'duncab7': str(dfca[26][1]), 'duncab8': str(dfca[27][1]),
        # Count:
        'ccea1': str(dfca[0][2])[:-2], 'ccea2': str(dfca[1][2])[:-2], 'ccea3': str(dfca[2][2])[:-2], 'ccea4': str(dfca[3][2])[:-2], 'ccea5': str(dfca[4][2])[:-2], 'ccea6': str(dfca[5][2])[:-2], 'ccea7': str(dfca[6][2])[:-2], 'ccea8': str(dfca[7][2])[:-2], 'ccea9': str(dfca[8][2])[:-2], 'ccea0': str(dfca[9][2])[:-2],
        'cceaa1': str(dfca[10][2])[:-2], 'cceaa2': str(dfca[11][2])[:-2], 'cceaa3': str(dfca[12][2])[:-2], 'cceaa4': str(dfca[13][2])[:-2], 'cceaa5': str(dfca[14][2])[:-2], 'cceaa6': str(dfca[15][2])[:-2], 'cceaa7': str(dfca[16][2])[:-2], 'cceaa8': str(dfca[17][2])[:-2], 'cceaa9': str(dfca[18][2])[:-2], 'cceaa0': str(dfca[19][2])[:-2],
        'cceab1': str(dfca[20][2])[:-2], 'cceab2': str(dfca[21][2])[:-2], 'cceab3': str(dfca[22][2])[:-2], 'cceab4': str(dfca[23][2])[:-2], 'cceab5': str(dfca[24][2])[:-2], 'cceab6': str(dfca[25][2])[:-2], 'cceab7': str(dfca[26][2])[:-2], 'cceab8': str(dfca[27][2])[:-2],
    # Enable User
        # Username:
        'eacu1': str(dfea[0][0]), 'eacu2': str(dfea[1][0]), 'eacu3': str(dfea[2][0]), 'eacu4': str(dfea[3][0]), 'eacu5': str(dfea[4][0]), 'eacu6': str(dfea[5][0]), 'eacu7': str(dfea[6][0]), 'eacu8': str(dfea[7][0]), 'eacu9': str(dfea[8][0]), 'eacu0': str(dfea[9][0]), 
        'eacua1': str(dfea[10][0]), 'eacua2': str(dfea[11][0]), 'eacua3': str(dfea[12][0]), 'eacua4': str(dfea[13][0]), 'eacua5': str(dfea[14][0]), 'eacua6': str(dfea[15][0]), 'eacua7': str(dfea[16][0]), 'eacua8': str(dfea[17][0]), 'eacua9': str(dfea[18][0]), 'eacua0': str(dfea[19][0]), 
        'eacub1': str(dfea[20][0]), 'eacub2': str(dfea[21][0]), 'eacub3': str(dfea[22][0]), 'eacub4': str(dfea[23][0]), 'eacub5': str(dfea[24][0]), 'eacub6': str(dfea[25][0]), 'eacub7': str(dfea[26][0]), 'eacub8': str(dfea[27][0]),
        # Destination User Name:
        'dunea1': str(dfea[0][1]), 'dunea2': str(dfea[1][1]), 'dunea3': str(dfea[2][1]), 'dunea4': str(dfea[3][1]), 'dunea5': str(dfea[4][1]), 'dunea6': str(dfea[5][1]), 'dunea7': str(dfea[6][1]), 'dunea8': str(dfea[7][1]), 'dunea9': str(dfea[8][1]), 'dunea0': str(dfea[9][1]), 
        'duneaa1': str(dfea[10][1]), 'duneaa2': str(dfea[11][1]), 'duneaa3': str(dfea[12][1]), 'duneaa4': str(dfea[13][1]), 'duneaa5': str(dfea[14][1]), 'duneaa6': str(dfea[15][1]), 'duneaa7': str(dfea[16][1]), 'duneaa8': str(dfea[17][1]), 'duneaa9': str(dfea[18][1]), 'duneaa0': str(dfea[19][1]), 
        'duneab1': str(dfea[20][1]), 'duneab2': str(dfea[21][1]), 'duneab3': str(dfea[22][1]), 'duneab4': str(dfea[23][1]), 'duneab5': str(dfea[24][1]), 'duneab6': str(dfea[25][1]), 'duneab7': str(dfea[26][1]), 'duneab8': str(dfea[27][1]),
        # Count:
        'ceea1': str(dfea[0][2])[:-2], 'ceea2': str(dfea[1][2])[:-2], 'ceea3': str(dfea[2][2])[:-2], 'ceea4': str(dfea[3][2])[:-2], 'ceea5': str(dfea[4][2])[:-2], 'ceea6': str(dfea[5][2])[:-2], 'ceea7': str(dfea[6][2])[:-2], 'ceea8': str(dfea[7][2])[:-2], 'ceea9': str(dfea[8][2])[:-2], 'ceea0': str(dfea[9][2])[:-2], 
        'ceeaa1': str(dfea[10][2])[:-2], 'ceeaa2': str(dfea[11][2])[:-2], 'ceeaa3': str(dfea[12][2])[:-2], 'ceeaa4': str(dfea[13][2])[:-2], 'ceeaa5': str(dfea[14][2])[:-2], 'ceeaa6': str(dfea[15][2])[:-2], 'ceeaa7': str(dfea[16][2])[:-2], 'ceeaa8': str(dfea[17][2])[:-2], 'ceeaa9': str(dfea[18][2])[:-2], 'ceeaa0': str(dfea[19][2])[:-2], 
        'ceeab1': str(dfea[20][2])[:-2], 'ceeab2': str(dfea[21][2])[:-2], 'ceeab3': str(dfea[22][2])[:-2], 'ceeab4': str(dfea[23][2])[:-2], 'ceeab5': str(dfea[24][2])[:-2], 'ceeab6': str(dfea[25][2])[:-2], 'ceeab7': str(dfea[26][2])[:-2], 'ceeab8': str(dfea[27][2])[:-2],
    # Reset Password
        # Username:
        'rpu1': str(dfrp[0][0]), 'rpu2': str(dfrp[1][0]), 'rpu3': str(dfrp[2][0]), 'rpu4': str(dfrp[3][0]), 'rpu5': str(dfrp[4][0]), 'rpu6': str(dfrp[5][0]), 'rpu7': str(dfrp[6][0]), 'rpu8': str(dfrp[7][0]), 'rpu9': str(dfrp[8][0]), 'rpu0': str(dfrp[9][0]), 
        # Destination Username:
        'rpun1': str(dfrp[0][1]), 'rpun2': str(dfrp[1][1]), 'rpun3': str(dfrp[2][1]), 'rpun4': str(dfrp[3][1]), 'rpun5': str(dfrp[4][1]), 'rpun6': str(dfrp[5][1]), 'rpun7': str(dfrp[6][1]), 'rpun8': str(dfrp[7][1]), 'rpun9': str(dfrp[8][1]), 'rpun0': str(dfrp[9][1]), 
        # Event Count:
        'rpc1': str(dfrp[0][2])[:-2], 'rpc2': str(dfrp[1][2])[:-2], 'rpc3': str(dfrp[2][2])[:-2], 'rpc4': str(dfrp[3][2])[:-2], 'rpc5': str(dfrp[4][2])[:-2], 'rpc6': str(dfrp[5][2])[:-2], 'rpc7': str(dfrp[6][2])[:-2], 'rpc8': str(dfrp[7][2])[:-2], 'rpc9': str(dfrp[8][2])[:-2], 'rpc0': str(dfrp[9][2])[:-2], 
    # Disable User
        # Username:
        'dau1': str(dfda[0][0]), 'dau2': str(dfda[1][0]), 'dau3': str(dfda[2][0]), 'dau4': str(dfda[3][0]), 'dau5': str(dfda[4][0]), 'dau6': str(dfda[5][0]), 'dau7': str(dfda[6][0]), 'dau8': str(dfda[7][0]), 'dau9': str(dfda[8][0]), 'dau0': str(dfda[9][0]), 
        # Destination Username:
        'daun1': str(dfda[0][1]), 'daun2': str(dfda[1][1]), 'daun3': str(dfda[2][1]), 'daun4': str(dfda[3][1]), 'daun5': str(dfda[4][1]), 'daun6': str(dfda[5][1]), 'daun7': str(dfda[6][1]), 'daun8': str(dfda[7][1]), 'daun9': str(dfda[8][1]), 'daun0': str(dfda[9][1]), 
        # Event Count:
        'dac1': str(dfda[0][2]), 'dac2': str(dfda[1][2]), 'dac3': str(dfda[2][2]), 'dac4': str(dfda[3][2]), 'dac5': str(dfda[4][2]), 'dac6': str(dfda[5][2]), 'dac7': str(dfda[6][2]), 'dac8': str(dfda[7][2]), 'dac9': str(dfda[8][2]), 'dac0': str(dfda[9][2]), 
}

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 1:  # Placeholder
                    replace_text_in_shape(shape, replacements)
                elif shape.shape_type == 19:  # Table
                    replace_text_in_table(shape.table, replacements)
        
        output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        prs.save(output_path)
        messagebox.showinfo("Success", "PowerPoint file has been processed and saved")

    except Exception as e:
        messagebox.showerror("Error", str(e))

#---------------------------------PROCESS OIC-BACKUP DATA---------------------------------#

def process_data_backup():
    excel_path1 = Fortigate_excel_entry.get()
    excel_path2 = PaloAlto_excel_entry.get()
    pptx_path = BACKUP_pptx_entry.get()

    # Read data from Excel files
    df = xw.Book(excel_path1)
    adf = xw.Book(excel_path2)

    # Read data from PowerPoint file
    prs = Presentation(pptx_path)

    # Read data from sheets
    #----------------------------------FORTIGATE----------------------------------#
    tfbs = df.sheets['TalkersFromBackup'] # Talkers From Backup
    tfbds = df.sheets['TalkersFromBackup_Denied'] # Talkers From Backup Denied
    ttbs = df.sheets['TalkerToBackup'] # Talkers To Backup
    ttbds = df.sheets['TalkerToBackup_Denied'] # Talkers To Backup Denied
    #----------------------------------PALOALTO----------------------------------#
    atfbs = adf.sheets['TalkersFromBackup'] # Talkers From Backup
    atfbds = adf.sheets['TalkersFromBackup_Denied'] # Talkers From Backup Denied
    attbs = adf.sheets['TalkersToBackup'] # Talkers To Backup
    attbds = adf.sheets['TalkersToBackup_Denied'] # Talkers To Backup Denied

    # Iterate values in sheets
    #----------------------------------FORTIGATE----------------------------------#
    tfb = tfbs.range("A2:E11").value
    tfbd = tfbds.range("A2:E11").value
    ttb = ttbs.range("A2:E11").value
    ttbd = ttbds.range("A2:E11").value
    #----------------------------------PALOALTO----------------------------------#
    atfb = atfbs.range("A2:E11").value
    atfbd = atfbds.range("A2:E11").value
    attb = attbs.range("A2:E11").value
    attbd = attbds.range("A2:E11").value

    # Define replacements
    replacements = {
        #----------------------------------FORTIGATE----------------------------------#
    # Talkers From Backup Allow
        # Source IP:
        'tfbip1': str(tfb[0][0]), 'tfbip2': str(tfb[1][0]), 'tfbip3': str(tfb[2][0]), 'tfbip4': str(tfb[3][0]), 'tfbip5': str(tfb[4][0]), 'tfbip6': str(tfb[5][0]), 'tfbip7': str(tfb[6][0]), 'tfbip8': str(tfb[7][0]), 'tfbip9': str(tfb[8][0]), 'tfbip0': str(tfb[9][0]), 
        # Source IP Name:
        # Destination IP:
        'tfbdip1': str(tfb[0][1]), 'tfbdip2': str(tfb[1][1]), 'tfbdip3': str(tfb[2][1]), 'tfbdip4': str(tfb[3][1]), 'tfbdip5': str(tfb[4][1]), 'tfbdip6': str(tfb[5][1]), 'tfbdip7': str(tfb[6][1]), 'tfbdip8': str(tfb[7][1]), 'tfbdip9': str(tfb[8][1]), 'tfbdip0': str(tfb[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'tfbdp1': str(tfb[0][2])[:-2], 'tfbdp2': str(tfb[1][2])[:-2], 'tfbdp3': str(tfb[2][2])[:-2], 'tfbdp4': str(tfb[3][2])[:-2], 'tfbdp5': str(tfb[4][2])[:-2], 'tfbdp6': str(tfb[5][2])[:-2], 'tfbdp7': str(tfb[6][2])[:-2], 'tfbdp8': str(tfb[7][2])[:-2], 'tfbdp9': str(tfb[8][2])[:-2], 'tfbdp0': str(tfb[9][2])[:-2], 
        # Device Action:
        'tfbda1': str(tfb[0][3]), 'tfbda2': str(tfb[1][3]), 'tfbda3': str(tfb[2][3]), 'tfbda4': str(tfb[3][3]), 'tfbda5': str(tfb[4][3]), 'tfbda6': str(tfb[5][3]), 'tfbda7': str(tfb[6][3]), 'tfbda8': str(tfb[7][3]), 'tfbda9': str(tfb[8][3]), 'tfbda0': str(tfb[9][3]), 
        # Event Count:
        'tfbec1': str(tfb[0][4])[:-2], 'tfbec2': str(tfb[1][4])[:-2], 'tfbec3': str(tfb[2][4])[:-2], 'tfbec4': str(tfb[3][4])[:-2], 'tfbec5': str(tfb[4][4])[:-2], 'tfbec6': str(tfb[5][4])[:-2], 'tfbec7': str(tfb[6][4])[:-2], 'tfbec8': str(tfb[7][4])[:-2], 'tfbec9': str(tfb[8][4])[:-2], 'tfbec0': str(tfb[9][4])[:-2],
    # Talkers From Backup Denied
        # Source IP:
        'tfbipd1': str(tfbd[0][0]), 'tfbipd2': str(tfbd[1][0]), 'tfbipd3': str(tfbd[2][0]), 'tfbipd4': str(tfbd[3][0]), 'tfbipd5': str(tfbd[4][0]), 'tfbipd6': str(tfbd[5][0]), 'tfbipd7': str(tfbd[6][0]), 'tfbipd8': str(tfbd[7][0]), 'tfbipd9': str(tfbd[8][0]), 'tfbipd0': str(tfbd[9][0]), 
        # Source IP Name:
        # Destination IP:
        'tfbdipd1': str(tfbd[0][1]), 'tfbdipd2': str(tfbd[1][1]), 'tfbdipd3': str(tfbd[2][1]), 'tfbdipd4': str(tfbd[3][1]), 'tfbdipd5': str(tfbd[4][1]), 'tfbdipd6': str(tfbd[5][1]), 'tfbdipd7': str(tfbd[6][1]), 'tfbdipd8': str(tfbd[7][1]), 'tfbdipd9': str(tfbd[8][1]), 'tfbdipd0': str(tfbd[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'tfbdpd1': str(tfbd[0][2])[:-2], 'tfbdpd1': str(tfbd[1][2])[:-2], 'tfbdpd1': str(tfbd[2][2])[:-2], 'tfbdpd1': str(tfbd[3][2])[:-2], 'tfbdpd1': str(tfbd[4][2])[:-2], 'tfbdpd1': str(tfbd[5][2])[:-2], 'tfbdpd1': str(tfbd[6][2])[:-2], 'tfbdpd1': str(tfbd[7][2])[:-2], 'tfbdpd1': str(tfbd[8][2])[:-2], 'tfbdpd1': str(tfbd[9][2])[:-2], 
        # Device Action:
        'tfbdad1': str(tfbd[0][3]), 'tfbdad2': str(tfbd[1][3]), 'tfbdad3': str(tfbd[2][3]), 'tfbdad4': str(tfbd[3][3]), 'tfbdad5': str(tfbd[4][3]), 'tfbdad6': str(tfbd[5][3]), 'tfbdad7': str(tfbd[6][3]), 'tfbdad8': str(tfbd[7][3]), 'tfbdad9': str(tfbd[8][3]), 'tfbdad0': str(tfbd[9][3]), 
        # Event Count:
        'tfbecd1': str(tfbd[0][4])[:-2], 'tfbecd2': str(tfbd[1][4])[:-2], 'tfbecd3': str(tfbd[2][4])[:-2], 'tfbecd4': str(tfbd[3][4])[:-2], 'tfbecd5': str(tfbd[4][4])[:-2], 'tfbecd6': str(tfbd[5][4])[:-2], 'tfbecd7': str(tfbd[6][4])[:-2], 'tfbecd8': str(tfbd[7][4])[:-2], 'tfbecd9': str(tfbd[8][4])[:-2], 'tfbecd0': str(tfbd[9][4])[:-2], 
    # Talkers To Backup Allow
        # Source IP:
        'ttbip1': str(ttb[0][0]), 'ttbip2': str(ttb[1][0]), 'ttbip3': str(ttb[2][0]), 'ttbip4': str(ttb[3][0]), 'ttbip5': str(ttb[4][0]), 'ttbip6': str(ttb[5][0]), 'ttbip7': str(ttb[6][0]), 'ttbip8': str(ttb[7][0]), 'ttbip9': str(ttb[8][0]), 'ttbip0': str(ttb[9][0]), 
        # Source IP Name:
        # Destination IP:
        'ttbdip1': str(ttb[0][1]), 'ttbdip2': str(ttb[1][1]), 'ttbdip3': str(ttb[2][1]), 'ttbdip4': str(ttb[3][1]), 'ttbdip5': str(ttb[4][1]), 'ttbdip6': str(ttb[5][1]), 'ttbdip7': str(ttb[6][1]), 'ttbdip8': str(ttb[7][1]), 'ttbdip9': str(ttb[8][1]), 'ttbdip0': str(ttb[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'ttbdp1': str(ttb[0][2])[:-2], 'ttbdp2': str(ttb[1][2])[:-2], 'ttbdp3': str(ttb[2][2])[:-2], 'ttbdp4': str(ttb[3][2])[:-2], 'ttbdp5': str(ttb[4][2])[:-2], 'ttbdp6': str(ttb[5][2])[:-2], 'ttbdp7': str(ttb[6][2])[:-2], 'ttbdp8': str(ttb[7][2])[:-2], 'ttbdp9': str(ttb[8][2])[:-2], 'ttbdp0': str(ttb[9][2])[:-2], 
        # Device Action:
        'ttbda1': str(ttb[0][3]), 'ttbda2': str(ttb[1][3]), 'ttbda3': str(ttb[2][3]), 'ttbda4': str(ttb[3][3]), 'ttbda5': str(ttb[4][3]), 'ttbda6': str(ttb[5][3]), 'ttbda7': str(ttb[6][3]), 'ttbda8': str(ttb[7][3]), 'ttbda9': str(ttb[8][3]), 'ttbda0': str(ttb[9][3]), 
        # Event Count:
        'ttbec1': str(ttb[0][4])[:-2], 'ttbec2': str(ttb[1][4])[:-2], 'ttbec3': str(ttb[2][4])[:-2], 'ttbec4': str(tfb[3][4])[:-2], 'ttbec5': str(ttb[4][4])[:-2], 'ttbec6': str(ttb[5][4])[:-2], 'ttbec7': str(ttb[6][4])[:-2], 'ttbec8': str(ttb[7][4])[:-2], 'ttbec9': str(ttb[8][4])[:-2], 'ttbec0': str(ttb[9][4])[:-2], 
    # Talkers To Backup Denied
        # Source IP:
        'ttbipd1': str(ttbd[0][0]), 'ttbipd2': str(ttbd[1][0]), 'ttbipd3': str(ttbd[2][0]), 'ttbipd4': str(ttbd[3][0]), 'ttbipd5': str(ttbd[4][0]), 'ttbipd6': str(ttbd[5][0]), 'ttbipd7': str(ttbd[6][0]), 'ttbipd8': str(ttbd[7][0]), 'ttbipd9': str(ttbd[8][0]), 'ttbipd0': str(ttbd[9][0]), 
        # Source IP Name:
        # Destination IP:
        'ttbdipd1': str(ttbd[0][1]), 'ttbdipd2': str(ttbd[1][1]), 'ttbdipd3': str(ttbd[2][1]), 'ttbdipd4': str(ttbd[3][1]), 'ttbdipd5': str(ttbd[4][1]), 'ttbdipd6': str(ttbd[5][1]), 'ttbdipd7': str(ttbd[6][1]), 'ttbdipd8': str(ttbd[7][1]), 'ttbdipd9': str(ttbd[8][1]), 'ttbdipd0': str(ttbd[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'ttbdpd1': str(ttbd[0][2])[:-2], 'ttbdpd2': str(ttbd[1][2])[:-2], 'ttbdpd3': str(ttbd[2][2])[:-2], 'ttbdpd4': str(ttbd[3][2])[:-2], 'ttbdpd5': str(ttbd[4][2])[:-2], 'ttbdpd6': str(ttbd[5][2])[:-2], 'ttbdpd7': str(ttbd[6][2])[:-2], 'ttbdpd8': str(ttbd[7][2])[:-2], 'ttbdpd9': str(ttbd[8][2])[:-2], 'ttbdpd0': str(ttbd[9][2])[:-2], 
        # Device Action:
        'ttbdad1': str(ttbd[0][3]), 'ttbdad2': str(ttbd[1][3]), 'ttbdad3': str(ttbd[2][3]), 'ttbdad4': str(ttbd[3][3]), 'ttbdad5': str(ttbd[4][3]), 'ttbdad6': str(ttbd[5][3]), 'ttbdad7': str(ttbd[6][3]),'ttbdad8': str(ttbd[7][3]), 'ttbdad9': str(ttbd[8][3]), 'ttbdad0': str(ttbd[9][3]), 
        # Event Count:
        'ttbecd1': str(ttbd[0][4])[:-2], 'ttbecd2': str(ttbd[1][4])[:-2], 'ttbecd3': str(ttbd[2][4])[:-2], 'ttbecd4': str(ttbd[3][4])[:-2], 'ttbecd5': str(ttbd[4][4])[:-2], 'ttbecd6': str(ttbd[5][4])[:-2], 'ttbecd7': str(ttbd[6][4])[:-2], 'ttbecd8': str(ttbd[7][4])[:-2], 'ttbecd9': str(ttbd[8][4])[:-2], 'ttbecd0': str(ttbd[9][4])[:-2], 
#----------------------------------PALOALTO----------------------------------#
    # Talkers From Backup
        # Source IP:
        'tfbipa1': str(atfb[0][0]), 'tfbipa2': str(atfb[1][0]), 'tfbipa3': str(atfb[2][0]), 'tfbipa4': str(atfb[3][0]), 'tfbipa5': str(atfb[4][0]), 'tfbipa6': str(atfb[5][0]), 'tfbipa7': str(atfb[6][0]), 'tfbipa8': str(atfb[7][0]), 'tfbipa9': str(atfb[8][0]), 'tfbipa0': str(atfb[9][0]), 
        # Source IP Name:
        # Destination IP:
        'tfbdipa1': str(atfb[0][1]), 'tfbdipa2': str(atfb[1][1]), 'tfbdipa3': str(atfb[2][1]), 'tfbdipa4': str(atfb[3][1]), 'tfbdipa5': str(atfb[4][1]), 'tfbdipa6': str(atfb[5][1]), 'tfbdipa7': str(atfb[6][1]), 'tfbdipa8': str(atfb[7][1]), 'tfbdipa9': str(atfb[8][1]), 'tfbdipa0': str(atfb[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'tfbdpa1': str(atfb[0][2])[:-2], 'tfbdpa2': str(atfb[1][2])[:-2], 'tfbdpa3': str(atfb[2][2])[:-2], 'tfbdpa4': str(atfb[3][2])[:-2], 'tfbdpa5': str(atfb[4][2])[:-2], 'tfbdpa6': str(atfb[5][2])[:-2], 'tfbdpa7': str(atfb[6][2])[:-2], 'tfbdpa8': str(atfb[7][2])[:-2], 'tfbdpa9': str(atfb[8][2])[:-2], 'tfbdpa0': str(atfb[9][2])[:-2], 
        # Device Action:
        'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 'tfbdaa1': str(atfb[0][3]), 
        # Event Count:
        'tfbeca1': str(atfb[0][4])[:-2], 'tfbeca2': str(atfb[1][4])[:-2], 'tfbeca3': str(atfb[2][4])[:-2], 'tfbeca4': str(atfb[3][4])[:-2], 'tfbeca5': str(atfb[4][4])[:-2], 'tfbeca6': str(atfb[5][4])[:-2], 'tfbeca7': str(atfb[6][4])[:-2], 'tfbeca8': str(atfb[7][4])[:-2], 'tfbeca9': str(atfb[8][4])[:-2], 'tfbeca0': str(atfb[9][4])[:-2], 
    # Talkers From Backup Denied
        # Source IP:
        'tfbipad1': str(atfbd[0][0]), 'tfbipad2': str(atfbd[1][0]), 'tfbipad3': str(atfbd[2][0]), 'tfbipad4': str(atfbd[3][0]), 'tfbipad5': str(atfbd[4][0]), 'tfbipad6': str(atfbd[5][0]), 'tfbipad7': str(atfbd[6][0]), 'tfbipad8': str(atfbd[7][0]), 'tfbipad9': str(atfbd[8][0]), 'tfbipad0': str(atfbd[9][0]), 
        # Source IP Name:
        # Destination IP:
        'tfbdipad1': str(atfbd[0][1]), 'tfbdipad2': str(atfbd[1][1]), 'tfbdipad3': str(atfbd[2][1]), 'tfbdipad4': str(atfbd[3][1]), 'tfbdipad5': str(atfbd[4][1]), 'tfbdipad6': str(atfbd[5][1]), 'tfbdipad7': str(atfbd[6][1]), 'tfbdipad8': str(atfbd[7][1]), 'tfbdipad9': str(atfbd[8][1]), 'tfbdipad0': str(atfbd[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'tfbdpad1': str(atfbd[0][2])[:-2], 'tfbdpad2': str(atfbd[1][2])[:-2], 'tfbdpad3': str(atfbd[2][2])[:-2], 'tfbdpad4': str(atfbd[3][2])[:-2], 'tfbdpad5': str(atfbd[4][2])[:-2], 'tfbdpad6': str(atfbd[5][2])[:-2], 'tfbdpad7': str(atfbd[6][2])[:-2], 'tfbdpad8': str(atfbd[7][2])[:-2], 'tfbdpad9': str(atfbd[8][2])[:-2], 'tfbdpad0': str(atfbd[9][2])[:-2], 
        # Device Action:
        'tfbdaad1': str(atfbd[0][3]), 'tfbdaad2': str(atfbd[1][3]), 'tfbdaad3': str(atfbd[2][3]), 'tfbdaad4': str(atfbd[3][3]), 'tfbdaad5': str(atfbd[4][3]), 'tfbdaad6': str(atfbd[5][3]), 'tfbdaad7': str(atfbd[6][3]), 'tfbdaad8': str(atfbd[7][3]), 'tfbdaad9': str(atfbd[8][3]), 'tfbdaad0': str(atfbd[9][3]), 
        # Event Count:
        'tfbecad1': str(atfbd[0][4]), 'tfbecad1': str(atfbd[1][4]), 'tfbecad1': str(atfbd[2][4]), 'tfbecad1': str(atfbd[3][4]), 'tfbecad1': str(atfbd[4][4]), 'tfbecad1': str(atfbd[5][4]), 'tfbecad1': str(atfbd[6][4]), 'tfbecad1': str(atfbd[7][4]), 'tfbecad1': str(atfbd[8][4]), 'tfbecad1': str(atfbd[9][4]), 
    # Talkers To Backup
        # Source IP:
        'ttbipa1': str(attb[0][0]), 'ttbipa2': str(attb[1][0]), 'ttbipa3': str(attb[2][0]), 'ttbipa4': str(attb[3][0]), 'ttbipa5': str(attb[4][0]), 'ttbipa6': str(attb[5][0]), 'ttbipa7': str(attb[6][0]), 'ttbipa8': str(attb[7][0]), 'ttbipa9': str(attb[8][0]), 'ttbipa0': str(attb[9][0]), 
        # Source IP Name:
        # Destination IP:
        'ttbdipa1': str(attb[0][1]), 'ttbdipa2': str(attb[1][1]), 'ttbdipa3': str(attb[2][1]), 'ttbdipa4': str(attb[3][1]), 'ttbdipa5': str(attb[4][1]), 'ttbdipa6': str(attb[5][1]), 'ttbdipa7': str(attb[6][1]), 'ttbdipa8': str(attb[7][1]), 'ttbdipa9': str(attb[8][1]), 'ttbdipa0': str(attb[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'ttbdpa1': str(attb[0][2])[:-2], 'ttbdpa2': str(attb[1][2])[:-2], 'ttbdpa3': str(attb[2][2])[:-2], 'ttbdpa4': str(attb[3][2])[:-2], 'ttbdpa5': str(attb[4][2])[:-2], 'ttbdpa6': str(attb[5][2])[:-2], 'ttbdpa7': str(attb[6][2])[:-2], 'ttbdpa8': str(attb[7][2])[:-2], 'ttbdpa9': str(attb[8][2])[:-2], 'ttbdpa0': str(attb[9][2])[:-2], 
        # Device Action:
        'ttbdaa1': str(attb[0][3]), 'ttbdaa2': str(attb[1][3]), 'ttbdaa3': str(attb[2][3]), 'ttbdaa4': str(attb[3][3]), 'ttbdaa5': str(attb[4][3]), 'ttbdaa6': str(attb[5][3]), 'ttbdaa7': str(attb[6][3]), 'ttbdaa8': str(attb[7][3]), 'ttbdaa9': str(attb[8][3]), 'ttbdaa0': str(attb[9][3]), 
        # Event Count:
        'ttbeca1': str(attb[0][4])[:-2], 'ttbeca2': str(attb[1][4])[:-2], 'ttbeca3': str(attb[2][4])[:-2], 'ttbeca4': str(attb[3][4])[:-2], 'ttbeca5': str(attb[4][4])[:-2], 'ttbeca6': str(attb[5][4])[:-2], 'ttbeca7': str(attb[6][4])[:-2], 'ttbeca8': str(attb[7][4])[:-2], 'ttbeca9': str(attb[8][4])[:-2], 'ttbeca0': str(attb[9][4])[:-2], 
    # Talkers To Backup Denied
        # Source IP:
        'ttbipda1': str(attbd[0][0]), 'ttbipda2': str(attbd[1][0]), 'ttbipda3': str(attbd[2][0]), 'ttbipda4': str(attbd[3][0]), 'ttbipda5': str(attbd[4][0]), 'ttbipda6': str(attbd[5][0]), 'ttbipda7': str(attbd[6][0]), 'ttbipda8': str(attbd[7][0]), 'ttbipda9': str(attbd[8][0]), 'ttbipda0': str(attbd[9][0]), 
        # Scource IP Name:
        # Destination IP:
        'ttbdipda1': str(attbd[0][1]), 'ttbdipda2': str(attbd[1][1]), 'ttbdipda3': str(attbd[2][1]), 'ttbdipda4': str(attbd[3][1]), 'ttbdipda5': str(attbd[4][1]), 'ttbdipda6': str(attbd[5][1]), 'ttbdipda7': str(attbd[6][1]), 'ttbdipda8': str(attbd[7][1]), 'ttbdipda9': str(attbd[8][1]), 'ttbdipda0': str(attbd[9][1]), 
        # Destination IP Name:
        # Destination Port:
        'ttbdpda1': str(attbd[0][2])[:-2], 'ttbdpda2': str(attbd[1][2])[:-2], 'ttbdpda3': str(attbd[2][2])[:-2], 'ttbdpda4': str(attbd[3][2])[:-2], 'ttbdpda5': str(attbd[4][2])[:-2], 'ttbdpda6': str(attbd[5][2])[:-2], 'ttbdpda7': str(attbd[6][2])[:-2], 'ttbdpda8': str(attbd[7][2])[:-2], 'ttbdpda9': str(attbd[8][2])[:-2], 'ttbdpda0': str(attbd[9][2])[:-2], 
        # Device Action:
        'ttbdada1': str(attbd[0][3]), 'ttbdada2': str(attbd[1][3]), 'ttbdada3': str(attbd[2][3]), 'ttbdada4': str(attbd[3][3]), 'ttbdada5': str(attbd[4][3]), 'ttbdada6': str(attbd[5][3]), 'ttbdada7': str(attbd[6][3]), 'ttbdada8': str(attbd[7][3]), 'ttbdada9': str(attbd[8][3]), 'ttbdada0': str(attbd[9][3]), 
        # Event Count:
        'ttbecda1': str(attbd[0][4])[:-2], 'ttbecda2': str(attbd[1][4])[:-2], 'ttbecda3': str(attbd[2][4])[:-2], 'ttbecda4': str(attbd[3][4])[:-2], 'ttbecda5': str(attbd[4][4])[:-2], 'ttbecda6': str(attbd[5][4])[:-2], 'ttbecda7': str(attbd[6][4])[:-2], 'ttbecda8': str(attbd[7][4])[:-2], 'ttbecda9': str(attbd[8][4])[:-2], 'ttbecda0': str(attbd[9][4])[:-2], 

    }

    # Replace text in PowerPoint
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            if shape.has_table:
                replace_text_in_table(shape.table, replacements)

    # Save the modified PowerPoint presentation
    output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    prs.save(output_path)
    messagebox.showinfo("Success", "Data processed and PowerPoint updated successfully!")

#---------------------------------PROCESS OIC-FG DATA---------------------------------#
def process_data_fg():
    excel_file = FG_excel_entry.get()
    pptx_file = FG_pptx_entry.get()

    # Read data from Excel files
    df = xw.Book(excel_file)

    # Read data from PowerPoint file
    prs = Presentation(pptx_file)

    dfp = df.sheets[0] # Protocol
    dfd = df.sheets[1] # Data Leak
    dfr = df.sheets[2] # Remote Control

    # Iterate values in sheets
    p = dfp.range("A2:H11").value # Protocol
    d = dfd.range("A2:H11").value # Data Leak
    r = dfr.range("A2:H11").value # Remote Control

    replacements = {
        # Protocol
            # Source:
            'sp1': str(p[0][0]), 'sp2': str(p[1][0]), 'sp3': str(p[2][0]), 'sp4': str(p[3][0]), 'sp5': str(p[4][0]), 'sp6': str(p[5][0]), 'sp7': str(p[6][0]), 'sp8': str(p[7][0]), 'sp9': str(p[8][0]), 'sp0': str(p[9][0]), 
            # Destination:
            'dp1': str(p[0][1]), 'dp2': str(p[1][1]), 'dp3': str(p[2][1]), 'dp4': str(p[3][1]), 'dp5': str(p[4][1]), 'dp6': str(p[5][1]), 'dp7': str(p[6][1]), 'dp8': str(p[7][1]), 'dp9': str(p[8][1]), 'dp0': str(p[9][1]), 
            # Port:
            'pp1': str(p[0][2])[:-2], 'pp2': str(p[1][2])[:-2], 'pp3': str(p[2][2])[:-2], 'pp4': str(p[3][2])[:-2], 'pp5': str(p[4][2])[:-2], 'pp6': str(p[5][2])[:-2], 'pp7': str(p[6][2])[:-2], 'pp8': str(p[7][2])[:-2], 'pp9': str(p[8][2])[:-2], 'pp0': str(p[9][2])[:-2], 
            # Protocol:
            'ppr1': str(p[0][3]), 'ppr2': str(p[1][3]), 'ppr3': str(p[2][3]), 'ppr4': str(p[3][3]), 'ppr5': str(p[4][3]), 'ppr6': str(p[5][3]), 'ppr7': str(p[6][3]), 'ppr8': str(p[7][3]), 'ppr9': str(p[8][3]), 'ppr0': str(p[9][3]), 
            # Action:
            'ap1': str(p[0][4]), 'ap2': str(p[1][4]), 'ap3': str(p[2][4]), 'ap4': str(p[3][4]), 'ap5': str(p[4][4]), 'ap6': str(p[5][4]), 'ap7': str(p[6][4]), 'ap8': str(p[7][4]), 'ap9': str(p[8][4]), 'ap0': str(p[9][4]), 
            # Host:
            'hp1': str(p[0][5]), 'hp2': str(p[1][5]), 'hp3': str(p[2][5]), 'hp4': str(p[3][5]), 'hp5': str(p[4][5]), 'hp6': str(p[5][5]), 'hp7': str(p[6][5]), 'hp8': str(p[7][5]), 'hp9': str(p[8][5]), 'hp0': str(p[9][5]), 
            # Count: 
            'cp1': str(p[0][6])[:-2], 'cp2': str(p[1][6])[:-2], 'cp3': str(p[2][6])[:-2], 'cp4': str(p[3][6])[:-2], 'cp5': str(p[4][6])[:-2], 'cp6': str(p[5][6])[:-2], 'cp7': str(p[6][6])[:-2], 'cp8': str(p[7][6])[:-2], 'cp9': str(p[8][6])[:-2], 'cp0': str(p[9][6])[:-2], 
        # Data Leak Risk:
            # Source:
            'sd1': str(d[0][0]), 'sd2': str(d[1][0]), 'sd3': str(d[2][0]), 'sd4': str(d[3][0]), 'sd5': str(d[4][0]), 'sd6': str(d[5][0]), 'sd7': str(d[6][0]), 'sd8': str(d[7][0]), 'sd9': str(d[8][0]), 'sd0': str(d[9][0]), 
            # Destination:
            'dd1': str(d[0][1]), 'dd2': str(d[1][1]), 'dd3': str(d[2][1]), 'dd4': str(d[3][1]), 'dd5': str(d[4][1]), 'dd6': str(d[5][1]), 'dd7': str(d[6][1]), 'dd8': str(d[7][1]), 'dd9': str(d[8][1]), 'dd0': str(d[9][1]), 
            # Port:
            'dp1': str(d[0][2])[:-2], 'dp2': str(d[0][2])[:-2], 'dp3': str(d[0][2])[:-2], 'dp4': str(d[0][2])[:-2], 'dp5': str(d[0][2])[:-2], 'dp6': str(d[0][2])[:-2], 'dp7': str(d[0][2])[:-2], 'dp8': str(d[0][2])[:-2], 'dp9': str(d[0][2])[:-2], 'dp0': str(d[0][2])[:-2], 
            # Protocol:
            'drp1': str(d[0][3]), 'drp2': str(d[1][3]), 'drp3': str(d[2][3]), 'drp4': str(d[3][3]), 'drp5': str(d[4][3]), 'drp6': str(d[5][3]), 'drp7': str(d[6][3]), 'drp8': str(d[7][3]), 'drp9': str(d[8][3]), 'drp0': str(d[9][3]), 
            # Action:
            'ad1': str(d[0][4]), 'ad2': str(d[1][4]), 'ad3': str(d[2][4]), 'ad4': str(d[3][4]), 'ad5': str(d[4][4]), 'ad6': str(d[5][4]), 'ad7': str(d[6][4]), 'ad8': str(d[7][4]), 'ad9': str(d[8][4]), 'ad0': str(d[9][4]), 
            # Host:
            'hd1': str(d[0][5]), 'hd2': str(d[1][5]), 'hd3': str(d[2][5]), 'hd4': str(d[3][5]), 'hd5': str(d[4][5]), 'hd6': str(d[5][5]), 'hd7': str(d[6][5]), 'hd8': str(d[7][5]), 'hd9': str(d[8][5]), 'hd0': str(d[9][5]), 
            # Count:
            'cd1': str(d[0][6])[:-2], 'cd2': str(d[1][6])[:-2], 'cd3': str(d[2][6])[:-2], 'cd4': str(d[3][6])[:-2], 'cd5': str(d[4][6])[:-2], 'cd6': str(d[5][6])[:-2], 'cd7': str(d[6][6])[:-2], 'cd8': str(d[7][6])[:-2], 'cd9': str(d[8][6])[:-2], 'cd0': str(d[9][6])[:-2], 
        # Remote Control Risk:
            # Source:
            'rp1': str(r[0][0]), 'rp2': str(r[1][0]), 'rp3': str(r[2][0]), 'rp4': str(r[3][0]), 'rp5': str(r[4][0]), 'rp6': str(r[5][0]), 'rp7': str(r[6][0]), 'rp8': str(r[7][0]), 'rp9': str(r[8][0]), 'rp0': str(r[9][0]), 
            # Destination:
            'rd1': str(d[0][1]), 'rd2': str(d[1][1]), 'rd3': str(d[2][1]), 'rd4': str(d[3][1]), 'rd5': str(d[4][1]), 'rd6': str(d[5][1]), 'rd7': str(d[6][1]), 'rd8': str(d[7][1]), 'rd9': str(d[9][1]), 'rd0': str(d[0][1]), 
            # Port:
            'pr1': str(d[0][2])[:-2], 'pr2': str(d[1][2])[:-2], 'pr3': str(d[2][2])[:-2], 'pr4': str(d[3][2])[:-2], 'pr5': str(d[4][2])[:-2], 'pr6': str(d[5][2])[:-2], 'pr7': str(d[6][2])[:-2], 'pr8': str(d[7][2])[:-2], 'pr9': str(d[8][2])[:-2], 'pr0': str(d[9][2])[:-2], 
            # Protocol:
            'prr1': str(d[0][3]), 'prr2': str(d[1][3]), 'prr3': str(d[2][3]), 'prr4': str(d[3][3]), 'prr5': str(d[4][3]), 'prr6': str(d[5][3]), 'prr7': str(d[6][3]), 'prr8': str(d[7][3]), 'prr9': str(d[8][3]), 'prr0': str(d[9][3]), 
            # Action:
            'ar1': str(d[0][4]), 'ar2': str(d[1][4]), 'ar3': str(d[2][4]), 'ar4': str(d[3][4]), 'ar5': str(d[4][4]), 'ar6': str(d[5][4]), 'ar7': str(d[6][4]), 'ar8': str(d[7][4]), 'ar9': str(d[8][4]), 'ar0': str(d[9][4]), 
            # Host:
            'hr1': str(d[0][5]), 'hr2': str(d[1][5]), 'hr3': str(d[2][5]), 'hr4': str(d[3][5]), 'hr5': str(d[4][5]), 'hr6': str(d[5][5]), 'hr7': str(d[6][5]), 'hr8': str(d[7][5]), 'hr9': str(d[8][5]), 'hr0': str(d[9][5]), 
            # Count:
            'cr1': str(d[0][6])[:-2], 'cr2': str(d[1][6])[:-2], 'cr3': str(d[2][6])[:-2], 'cr4': str(d[3][6])[:-2], 'cr5': str(d[4][6])[:-2], 'cr6': str(d[5][6])[:-2], 'cr7': str(d[6][6])[:-2], 'cr8': str(d[7][6])[:-2], 'cr9': str(d[8][6])[:-2], 'cr0': str(d[9][6])[:-2], 
    }

    # Replace text in PowerPoint
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            if shape.has_table:
                replace_text_in_table(shape.table, replacements)

    # Save the modified PowerPoint presentation
    output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    prs.save(output_path)
    messagebox.showinfo("Success", "Data processed and PowerPoint updated successfully!")

#---------------------------------PROCESS OIC-INCAPSULA DATA---------------------------------#
def process_data_incapsula():
    excel_path = INCAPSULA_excel_entry.get()
    pptx_path = INCAPSULA_pptx_entry.get()
        
    if not excel_path or not pptx_path:
        messagebox.showerror("Error", "Please select both Excel and PowerPoint files")
        return

    # Read data from Excel
    df = xw.Book(excel_path)
        
    # Read data from pptx
    prs = Presentation(pptx_path)

    # Read data from sheets
    dfass = df.sheets[0] # Authentication Successful

    # Iterate values in sheets
    dfas = dfass.range("A2:E11").value # Authentication Successful

    replacements = {
    # Incapsula Details
        # Action:
        'a1': str(dfas[0][1]), 'a2': str(dfas[1][1]), 'a3': str(dfas[2][1]), 'a4': str(dfas[3][1]), 'a5': str(dfas[4][1]), 'a6': str(dfas[5][1]), 'a7': str(dfas[6][1]), 'a8': str(dfas[7][1]), 'a9': str(dfas[8][1]), 'a0': str(dfas[9][1]), 
        # Source Address:
        's1': str(dfas[0][2]), 's2': str(dfas[1][2]), 's3': str(dfas[2][2]), 's4': str(dfas[3][2]), 's5': str(dfas[4][2]), 's6': str(dfas[5][2]), 's7': str(dfas[6][2]), 's8': str(dfas[7][2]), 's9': str(dfas[8][2]), 's0': str(dfas[9][2]), 
        # Country:
        # Destination:
        'd1': str(dfas[0][3]), 'd2': str(dfas[1][3]), 'd3': str(dfas[2][3]), 'd4': str(dfas[3][3]), 'd5': str(dfas[4][3]), 'd6': str(dfas[5][3]), 'd7': str(dfas[6][3]), 'd8': str(dfas[7][3]), 'd9': str(dfas[8][3]), 'd0': str(dfas[9][3]), 
        # Event Count:
        'e1': str(dfas[0][4])[:-2], 'e2': str(dfas[1][4])[:-2], 'e3': str(dfas[2][4])[:-2], 'e4': str(dfas[3][4])[:-2], 'e5': str(dfas[4][4])[:-2], 'e6': str(dfas[5][4])[:-2], 'e7': str(dfas[6][4])[:-2], 'e8': str(dfas[7][4])[:-2], 'e9': str(dfas[8][4])[:-2], 'e0': str(dfas[9][4])[:-2], 
}
    
    # Replace text in PowerPoint
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            if shape.has_table:
                replace_text_in_table(shape.table, replacements)

    # Save the modified PowerPoint presentation
    output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    prs.save(output_path)
    messagebox.showinfo("Success", "Data processed and PowerPoint updated successfully!")

#---------------------------------PROCESS OIC-PALOALTO DATA---------------------------------#
def process_data_paloalto():
    excel_file = PA_excel_entry.get()
    pptx_file = PA_pptx_entry.get()

    # Read data from Excel files
    df = xw.Book(excel_file)

    # Read data from PowerPoint file
    prs = Presentation(pptx_file)

    # Read data from sheets
    dfpib = df.sheets['Protocol_Inbound'] # Protocol Inbound
    dfpob = df.sheets['Protocol_Outbound'] # Protocol Outbound
    dflri = df.sheets['Data_Leak_Risk_Inbound'] # Data Leak Risk Inbound
    dflro = df.sheets['Data_Leak_Risk_Outbound'] # Data Leak Risk Outbound
    dfrib = df.sheets['Remote_Control_Inbound'] # Remote Control Outbound
    dfrob = df.sheets['Remote_Control_Outbound'] # Remote Countrol Outbound


    # Iterate values in sheets
    pib = dfpib.range("A2:H6").value # Protocol Inbound
    pob = dfpob.range("A2:H6").value # Protocol Outbound
    lri = dflri.range("A2:H6").value # Data Leak Risk Inbound
    lro = dflro.range("A2:H6").value # Data Leak RIsk Outbound
    rib = dfrib.range("A2:H6").value # Remote Inbound
    rob = dfrob.range("A2:H6").value # Remote Outbound

    replacements = {
    # Protocol Inbound
        # Source IP:
        'sib1': str(pib[0][0]), 'sib2': str(pib[1][0]), 'sib3': str(pib[2][0]), 'sib4': str(pib[3][0]), 'sib5': str(pib[4][0]), 
        # Destination:
        'dib1': str(pib[0][2]), 'dib2': str(pib[1][2]), 'dib3': str(pib[2][2]), 'dib4': str(pib[3][2]), 'dib5': str(pib[4][2]), 
        # Country
        'cib1': str(pib[0][1]), 'cib2': str(pib[1][1]), 'cib3': str(pib[2][1]), 'cib4': str(pib[3][1]), 'cib5': str(pib[4][1]), 
        # Port:
        'pib1': str(pib[0][3])[:-2], 'pib2': str(pib[1][3])[:-2], 'pib3': str(pib[2][3])[:-2], 'pib4': str(pib[3][3])[:-2], 'pib5': str(pib[4][3])[:-2], 
        # Protocol:
        'prib1': str(pib[0][4]), 'prib2': str(pib[1][4]), 'prib3': str(pib[2][4]), 'prib4': str(pib[3][4]), 'prib5': str(pib[4][4]), 
        # Action:
        'aib1': str(pib[0][5]), 'aib2': str(pib[1][5]), 'aib3': str(pib[2][5]), 'aib4': str(pib[3][5]), 'aib5': str(pib[4][5]), 
        # Host:
        'hib1': str(pib[0][6]), 'hib2': str(pib[1][6]), 'hib3': str(pib[2][6]), 'hib4': str(pib[3][6]), 'hib5': str(pib[4][6]), 
        # Count:
        'icb1': str(pib[0][7])[:-2], 'icb2': str(pib[1][7])[:-2], 'icb3': str(pib[2][7])[:-2], 'icb4': str(pib[3][7])[:-2], 'icb5': str(pib[4][7])[:-2], 
    # Protocol Outbound
        # Source IP:
        'sob1': str(pob[0][0]), 'sob2': str(pob[1][0]), 'sob3': str(pob[2][0]), 'sob4': str(pob[3][0]), 'sob5': str(pob[4][0]),
        # Destination:
        'dob1': str(pob[0][1]), 'dob2': str(pob[1][1]), 'dob3': str(pob[2][1]), 'dob4': str(pob[3][1]), 'dob5': str(pob[4][1]),
        # Country
        'cob1': str(pob[0][2]), 'cob2': str(pob[1][2]), 'cob3': str(pob[2][2]), 'cob4': str(pob[3][2]), 'cob5': str(pob[4][2]), 
        # Port
        'pob1': str(pob[0][3])[:-2], 'pob2': str(pob[1][3])[:-2], 'pob3': str(pob[2][3])[:-2], 'pob4': str(pob[3][3])[:-2], 'pob5': str(pob[4][3])[:-2],
        # Protocol
        'prob1': str(pob[0][4]), 'prob2': str(pob[1][4]), 'prob3': str(pob[2][4]), 'prob4': str(pob[3][4]), 'prob5': str(pob[4][4]),
        # Action
        'aob1': str(pob[0][5]), 'aob2': str(pob[1][5]), 'aob3': str(pob[2][5]), 'aob4': str(pob[3][5]), 'aob5': str(pob[4][5]),
        # Host
        'hob1': str(pob[0][6]), 'hob2': str(pob[1][6]), 'hob3': str(pob[2][6]), 'hob4': str(pob[3][6]), 'hob5': str(pob[4][6]), 
        # Count
        'ocb1': str(pob[0][7])[:-2], 'ocb2': str(pob[1][7])[:-2], 'ocb3': str(pob[2][7])[:-2], 'ocb4': str(pob[3][7])[:-2], 'ocb5': str(pob[4][7])[:-2], 
    # Data Leak Inbound
        # Source:
        'sibd1': str(lri[0][0]), 'sibd2': str(lri[1][0]), 'sibd3': str(lri[2][0]), 'sibd4': str(lri[3][0]), 'sibd5': str(lri[4][0]),
        # Destination:
        'dibd1': str(lri[0][2]), 'dibd2': str(lri[1][2]), 'dibd3': str(lri[2][2]), 'dibd4': str(lri[3][2]), 'dibd5': str(lri[4][2]),  
        # Country:
        'cibds1': str(lri[0][1]), 'cibds2': str(lri[1][1]), 'cibds3': str(lri[2][1]), 'cibds4': str(lri[3][1]), 'cibds5': str(lri[4][1]), 
        # Port:
        'pibd1': str(lri[0][3])[:-2], 'pibd2': str(lri[1][3])[:-2], 'pibd3': str(lri[2][3])[:-2], 'pibd4': str(lri[3][3])[:-2], 'pibd5': str(lri[4][3])[:-2], 
        # Protocol:
        'pribd1': str(lri[0][4]), 'pribd2': str(lri[1][4]), 'pribd3': str(lri[2][4]), 'pribd4': str(lri[3][4]), 'pribd5': str(lri[4][4]), 
        # Action:
        'aibd1': str(lri[0][5]), 'aibd2': str(lri[1][5]), 'aibd3': str(lri[2][5]), 'aibd4': str(lri[3][5]), 'aibd5': str(lri[4][5]), 
        # Host:
        'hibd1': str(lri[0][6]), 'hibd2': str(lri[1][6]), 'hibd3': str(lri[2][6]), 'hibd4': str(lri[3][6]), 'hibd5': str(lri[4][6]), 
        # Count:
        'icbd1': str(lri[0][7])[:-2], 'icbd1': str(lri[0][7])[:-2], 'icbd1': str(lri[0][7])[:-2], 'icbd1': str(lri[0][7])[:-2], 'icbd1': str(lri[0][7])[:-2], 
    # Data Leak Outbound
        # Source:
        'sobd1': str(lro[0][0]), 'sobd2': str(lro[1][0]), 'sobd3': str(lro[2][0]), 'sobd4': str(lro[3][0]), 'sobd5': str(lro[4][0]), 
        # Destination:
        'dobd1': str(lro[0][1]), 'dobd2': str(lro[1][1]), 'dobd3': str(lro[2][1]), 'dobd4': str(lro[3][1]), 'dobd5': str(lro[4][1]), 
        # Country:
        'cobds1': str(lro[0][2]), 'cobds2': str(lro[1][2]), 'cobds3': str(lro[2][2]), 'cobds4': str(lro[3][2]), 'cobds5': str(lro[4][2]), 
        # Port:
        'pobd1': str(lro[0][3])[:-2], 'pobd2': str(lro[1][3])[:-2], 'pobd3': str(lro[2][3])[:-2], 'pobd4': str(lro[3][3])[:-2], 'pobd5': str(lro[4][3])[:-2], 
        # Protocol:
        'probd1': str(lro[0][4]), 'probd2': str(lro[1][4]), 'probd3': str(lro[2][4]), 'probd4': str(lro[3][4]), 'probd5': str(lro[4][4]), 
        # Action:
        'aobd1': str(lro[0][5]), 'aobd2': str(lro[1][5]), 'aobd3': str(lro[2][5]), 'aobd4': str(lro[3][5]), 'aobd5': str(lro[4][5]), 
        # Host:
        'hobd1': str(lro[0][6]), 'hobd2': str(lro[1][6]), 'hobd3': str(lro[2][6]), 'hobd4': str(lro[3][6]), 'hobd5': str(lro[4][6]), 
        # Count:
        'ocbd1': str(lro[0][7])[:-2], 'ocbd2': str(lro[1][7])[:-2], 'ocbd3': str(lro[2][7])[:-2], 'ocbd4': str(lro[3][7])[:-2], 'ocbd5': str(lro[4][7])[:-2],
    # Remote Control Risk Inbound
        # Source:
        'sibr1': str(rib[0][0]), 'sibr2': str(rib[1][0]), 'sibr3': str(rib[2][0]), 'sibr4': str(rib[3][0]), 'sibr5': str(rib[4][0]), 
        # Destination: 
        'dibr1': str(rib[0][2]), 'dibr2': str(rib[1][2]), 'dibr3': str(rib[2][2]), 'dibr4': str(rib[3][2]), 'dibr5': str(rib[4][2]), 
        # Country:
        'cibr1': str(rib[0][1]), 'cibr2': str(rib[1][1]), 'cibr3': str(rib[2][1]), 'cibr4': str(rib[3][1]), 'cibr5': str(rib[4][1]), 
        # Port:
        'pibr1': str(rib[0][3])[:-2], 'pibr2': str(rib[1][3])[:-2], 'pibr3': str(rib[2][3])[:-2], 'pibr4': str(rib[3][3])[:-2], 'pibr5': str(rib[4][3])[:-2], 
        # Protocol:
        'pribr1': str(rib[0][4]), 'pribr2': str(rib[1][4]), 'pribr3': str(rib[2][4]), 'pribr4': str(rib[3][4]), 'pribr5': str(rib[4][4]), 
        # Action:
        'aibr1': str(rib[0][5]), 'aibr2': str(rib[1][5]), 'aibr3': str(rib[2][5]), 'aibr4': str(rib[3][5]), 'aibr5': str(rib[4][5]), 
        # Host:
        'hibr1': str(rib[0][6]), 'hibr2': str(rib[1][6]), 'hibr3': str(rib[2][6]), 'hibr4': str(rib[3][6]), 'hibr5': str(rib[4][6]), 
        # Count:
        'icbr1': str(rib[0][7])[:-2], 'icbr2': str(rib[1][7])[:-2], 'icbr3': str(rib[2][7])[:-2], 'icbr4': str(rib[3][7])[:-2], 'icbr5': str(rib[4][7])[:-2],
    # Remote Control Risk Outbound 
        # Source:
        'sobr1': str(rob[0][0]), 'sobr2': str(rob[1][0]), 'sobr3': str(rob[2][0]), 'sobr4': str(rob[3][0]), 'sobr5': str(rob[4][0]), 
        # Destination:
        'dobr1': str(rob[0][1]), 'dobr2': str(rob[1][1]), 'dobr3': str(rob[2][1]), 'dobr4': str(rob[3][1]), 'dobr5': str(rob[4][1]), 
        # Country:
        'cobr1': str(rob[0][2]), 'cobr2': str(rob[1][2]), 'cobr3': str(rob[2][2]), 'cobr4': str(rob[3][2]), 'cobr5': str(rob[4][2]), 
        # Port:
        'pobr1': str(rob[0][3])[:-2], 'pobr2': str(rob[1][3])[:-2], 'pobr3': str(rob[2][3])[:-2], 'pobr4': str(rob[3][3])[:-2], 'pobr5': str(rob[4][3])[:-2], 
        # Protocol:
        'probr1': str(rob[0][4]), 'probr2': str(rob[1][4]), 'probr3': str(rob[2][4]), 'probr4': str(rob[3][4]), 'probr5': str(rob[4][4]),  
        # Action:
        'aobr1': str(rob[0][5]), 'aobr2': str(rob[1][5]), 'aobr3': str(rob[2][5]), 'aobr4': str(rob[3][5]), 'aobr5': str(rob[4][5]), 
        # Host:
        'hobr1': str(rob[0][6]), 'hobr2': str(rob[1][6]), 'hobr3': str(rob[2][6]), 'hobr4': str(rob[3][6]), 'hobr5': str(rob[4][6]), 
        # Count:
        'ocbr1': str(rob[0][7])[:-2], 'ocbr2': str(rob[1][7])[:-2], 'ocbr3': str(rob[2][7])[:-2], 'ocbr4': str(rob[3][7])[:-2], 'ocbr5': str(rob[4][7])[:-2], 
        
}   
     # Replace text in PowerPoint
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            if shape.has_table:
                replace_text_in_table(shape.table, replacements)
    
    # Save the modified PowerPoint presentation
    output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    prs.save(output_path)
    messagebox.showinfo("Success", "Data processed and PowerPoint updated successfully!")

#---------------------------------PROCESS OIC-TRENDMICRO DATA---------------------------------#
def process_data_trendmicro():
    try:
        excel_file = TRENDMICRO_excel_entry.get()
        pptx_file = TRENDMICRO_pptx_entry.get()
        
        if not TRENDMICRO_excel_entry or not TRENDMICRO_pptx_entry:
            messagebox.showerror("Error", "Please select both Excel and PowerPoint files")
            return

        # Read data from Excel
        df = xw.Book(excel_file)
        
        # Read data from pptx
        prs = Presentation(pptx_file)

        # Read data from sheets
        dfas = df.sheets[0] # Vulnerability Action reset

        # Iterate values in sheets
        a = dfas.range("A2:F33").value # Vulnerability Action reset
        print(a)
        replacements = {
        # Vulnerability Action reset
        # Signature:
        's1': str(a[0][0]), 's2': str(a[1][0]), 's3': str(a[2][0]), 's4': str(a[3][0]), 's5': str(a[4][0]), 's6': str(a[5][0]), 's7': str(a[6][0]), 's8': str(a[7][0]), 's9': str(a[8][0]), 's0': str(a[9][0]), 
        'sz1': str(a[10][0]), 'sz2': str(a[11][0]), 'sz3': str(a[12][0]), 'sz4': str(a[13][0]), 'sz5': str(a[14][0]), 'sz6': str(a[15][0]), 'sz7': str(a[16][0]), 'sz8': str(a[17][0]), 
        # Source Address:
        'sa1': str(a[0][1]), 'sa2': str(a[1][1]), 'sa3': str(a[2][1]), 'sa4': str(a[3][1]), 'sa5': str(a[4][1]), 'sa6': str(a[5][1]), 'sa7': str(a[6][1]), 'sa8': str(a[7][1]), 'sa9': str(a[8][1]), 'sa0': str(a[9][1]), 
        'saa1': str(a[10][1]), 'saa2': str(a[11][1]), 'saa3': str(a[12][1]), 'saa4': str(a[13][1]), 'saa5': str(a[14][1]), 'saa6': str(a[15][1]), 'saa7': str(a[16][1]), 'saa8': str(a[17][1]),  
        # Host Name:
        # Destination Port:
        'd1': str(a[0][2]), 'd2': str(a[1][2]), 'd3': str(a[2][2]), 'd4': str(a[3][2]), 'd5': str(a[4][2]), 'd6': str(a[5][2]), 'd7': str(a[6][2]), 'd8': str(a[7][2]), 'd9': str(a[8][2]), 'd0': str(a[9][2]), 
        'da1': str(a[10][2]), 'da2': str(a[11][2]), 'da3': str(a[12][2]), 'da4': str(a[13][2]), 'da5': str(a[14][2]), 'da6': str(a[15][2]), 'da7': str(a[16][2]), 'da8': str(a[17][2]), 'da9': str(a[18][2]), 'da0': str(a[19][2]), 
        # Port:
        'p1': str(a[0][3])[:-2], 'p2': str(a[1][3])[:-2], 'p3': str(a[2][3])[:-2], 'p4': str(a[3][3])[:-2], 'p5': str(a[4][3])[:-2], 'p6': str(a[5][3])[:-2], 'p7': str(a[6][3])[:-2], 'p8': str(a[7][3])[:-2], 'p9': str(a[8][3])[:-2], 'p0': str(a[9][3])[:-2], 
        'pa1': str(a[10][3])[:-2], 'pa2': str(a[11][3])[:-2], 'pa3': str(a[12][3])[:-2], 'pa4': str(a[13][3])[:-2], 'pa5': str(a[14][3])[:-2], 'pa6': str(a[15][3])[:-2], 'pa7': str(a[16][3])[:-2], 'pa8': str(a[17][3])[:-2], 
        # Action:
        'ac1': str(a[0][4]), 'ac2': str(a[1][4]), 'ac3': str(a[2][4]), 'ac4': str(a[3][4]), 'ac5': str(a[4][4]), 'ac6': str(a[5][4]), 'ac7': str(a[6][4]), 'ac8': str(a[7][4]), 'ac9': str(a[8][4]), 'ac0': str(a[9][4]), 
        'aca1': str(a[10][4]), 'aca2': str(a[11][4]), 'aca3': str(a[12][4]), 'aca4': str(a[13][4]), 'aca5': str(a[14][4]), 'aca6': str(a[15][4]), 'aca7': str(a[16][4]), 'aca8': str(a[17][4]), 
        # Event Count:
        'ec1': str(a[0][5])[:-2], 'ec2': str(a[1][5])[:-2], 'ec3': str(a[2][5])[:-2], 'ec4': str(a[3][5])[:-2], 'ec5': str(a[4][5])[:-2], 'ec6': str(a[5][5])[:-2], 'ec7': str(a[6][5])[:-2], 'ec8': str(a[7][5])[:-2], 'ec9': str(a[8][5])[:-2], 'ec0': str(a[9][5])[:-2], 
        'eca1': str(a[10][5])[:-2], 'eca2': str(a[11][5])[:-2], 'eca3': str(a[12][5])[:-2], 'eca4': str(a[13][5])[:-2], 'eca5': str(a[14][5])[:-2], 'eca6': str(a[15][5])[:-2], 'eca7': str(a[16][5])[:-2], 'eca8': str(a[17][5])[:-2], 
        
}
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 1:  # Placeholder
                    replace_text_in_shape(shape, replacements)
                elif shape.shape_type == 19:  # Table
                    replace_text_in_table(shape.table, replacements)
        
        output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
        prs.save(output_path)
        messagebox.showinfo("Success", "PowerPoint file has been processed and saved")

    except Exception as e:
        messagebox.showerror("Error", str(e))

#---------------------------------PROCESS OIC-VPN DATA---------------------------------#
def process_data_vpn():
    excel_file = VPN_Success_excel_entry.get()
    excel_file2 = VPN_Failed_excel_entry.get()
    pptx_file = VPN_pptx_entry.get()

    # Read data from Excel files
    df = xw.Book(excel_file)
    dff = xw.Book(excel_file2)

    # Read data from PowerPoint file
    prs = Presentation(pptx_file)

    dfass = df.sheets[2] # Authentication Success
    dfafs = dff.sheets[2] # Authentication Failed 

    # Iterate values in sheets
    dfas = dfass.range("A2:C11").value # Authentication Successful
    dfaf = dfafs.range("A2:C11").value # Authentication Failure

    replacements = {
        # VPN Authentication Successful
            # IP Source Address:
            'visp1': str(dfas[0][0]), 'visp2': str(dfas[1][0]), 'visp3': str(dfas[2][0]), 'visp4': str(dfas[3][0]), 'visp5': str(dfas[4][0]), 'visp6': str(dfas[5][0]), 'visp7': str(dfas[6][0]), 'visp8': str(dfas[7][0]), 'visp9': str(dfas[8][0]), 'visp0': str(dfas[9][0]), 
            # Username:
            'visu1': str(dfas[0][1]), 'visu2': str(dfas[1][1]), 'visu3': str(dfas[2][1]), 'visu4': str(dfas[3][1]), 'visu5': str(dfas[4][1]), 'visu6': str(dfas[5][1]), 'visu7': str(dfas[6][1]), 'visu8': str(dfas[7][1]), 'visu9': str(dfas[8][1]), 'visu0': str(dfas[9][1]), 
            # Event Count:
            'viecs1': str(dfas[0][2])[:-2], 'viecs2': str(dfas[1][2])[:-2], 'viecs3': str(dfas[2][2])[:-2], 'viecs4': str(dfas[3][2])[:-2], 'viecs5': str(dfas[4][2])[:-2], 'viecs6': str(dfas[5][2])[:-2], 'viecs7': str(dfas[6][2])[:-2], 'viecs8': str(dfas[7][2])[:-2], 'viecs9': str(dfas[8][2])[:-2], 'viecs0': str(dfas[9][2])[:-2], 
        # VPN Authentication Failed
            # Username:
            'vifu1': str(dfaf[0][0]), 'vifu2': str(dfaf[1][0]), 'vifu3': str(dfaf[2][0]), 'vifu4': str(dfaf[3][0]), 'vifu5': str(dfaf[4][0]), 'vifu6': str(dfaf[5][0]), 'vifu7': str(dfaf[6][0]), 'vifu8': str(dfaf[7][0]), 'vifu9': str(dfaf[8][0]), 'vifu0': str(dfaf[9][0]), 
            # IP Source Address:
            'vifp1': str(dfaf[0][1]), 'vifp2': str(dfaf[1][1]), 'vifp3': str(dfaf[2][1]), 'vifp4': str(dfaf[3][1]), 'vifp5': str(dfaf[4][1]), 'vifp6': str(dfaf[5][1]), 'vifp7': str(dfaf[6][1]), 'vifp8': str(dfaf[7][1]), 'vifp9': str(dfaf[8][1]), 'vifp0': str(dfaf[9][1]), 
            # IP Source Address Name:
            # ISP/Country:
            # Event Count:
            'viecf1': str(dfaf[0][2])[:-2], 'viecf2': str(dfaf[1][2])[:-2], 'viecf3': str(dfaf[2][2])[:-2], 'viecf4': str(dfaf[3][2])[:-2], 'viecf5': str(dfaf[4][2])[:-2], 'viecf6': str(dfaf[5][2])[:-2], 'viecf7': str(dfaf[6][2])[:-2], 'viecf8': str(dfaf[7][2])[:-2], 'viecf9': str(dfaf[8][2])[:-2], 'viecf0': str(dfaf[9][2])[:-2], 
}

    # Replace text in PowerPoint
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text_in_shape(shape, replacements)
            if shape.has_table:
                replace_text_in_table(shape.table, replacements)

    # Save the modified PowerPoint presentation
    output_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
    prs.save(output_path)
    messagebox.showinfo("Success", "Data processed and PowerPoint updated successfully!")

#---------------------------------INTERFACE---------------------------------#
# Function to open file dialog
def open_file_dialog(entry):
    file_path = filedialog.askopenfilename()
    entry.delete(0, tk.END)
    entry.insert(0, file_path)

# Create main window
root = tk.Tk()
root.title("OIC Automation")

# Create and place labels and entry fields
#---------------------------------OIC-AD---------------------------------#
tk.Label(root, text="OIC-AD").grid(row=0, column=0, padx=10, pady=10)

tk.Label(root, text="AD Excel File:").grid(row=1, column=0, padx=10, pady=5)
AD_excel_entry = tk.Entry(root, width=50)
AD_excel_entry.grid(row=1, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(AD_excel_entry)).grid(row=1, column=2, padx=10, pady=5)

tk.Label(root, text="AD PowerPoint File:").grid(row=2, column=0, padx=10, pady=5)
AD_pptx_entry = tk.Entry(root, width=50)
AD_pptx_entry.grid(row=2, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(AD_pptx_entry)).grid(row=2, column=2, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_ad).grid(row=3, column=0, columnspan=3, pady=20)

#---------------------------------OIC-BACKUP---------------------------------#
tk.Label(root, text="OIC-BACKUP").grid(row=4, column=0, padx=10, pady=10)

tk.Label(root, text="Fortigate Excel File:").grid(row=5, column=0, padx=10, pady=5)
Fortigate_excel_entry = tk.Entry(root, width=50)
Fortigate_excel_entry.grid(row=5, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(Fortigate_excel_entry)).grid(row=5, column=2, padx=10, pady=5)

tk.Label(root, text="PaloAlto Excel File:").grid(row=6, column=0, padx=10, pady=5)
PaloAlto_excel_entry = tk.Entry(root, width=50)
PaloAlto_excel_entry.grid(row=6, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(PaloAlto_excel_entry)).grid(row=6, column=2, padx=10, pady=5)

tk.Label(root, text="PowerPoint File:").grid(row=7, column=0, padx=10, pady=5)
BACKUP_pptx_entry = tk.Entry(root, width=50)
BACKUP_pptx_entry.grid(row=7, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(BACKUP_pptx_entry)).grid(row=7, column=2, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_backup).grid(row=8, column=0, columnspan=3, padx=20)

#---------------------------------OIC-FG---------------------------------#
tk.Label(root, text="OIC-FG").grid(row=9, column=0, padx=10, pady=10)

tk.Label(root, text="FG Excel File:").grid(row=10, column=0, padx=10, pady=5)
FG_excel_entry = tk.Entry(root, width=50)
FG_excel_entry.grid(row=10, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(FG_excel_entry)).grid(row=10, column=2, padx=10, pady=5)

tk.Label(root, text="FG PPTX File:").grid(row=11, column=0, padx=10, pady=5)
FG_pptx_entry = tk.Entry(root, width=50)
FG_pptx_entry.grid(row=11, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(FG_pptx_entry)).grid(row=11, column=2, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_fg).grid(row=12, column=0, columnspan=3, pady=20)

#---------------------------------OIC-INCAPSULA---------------------------------#
tk.Label(root, text="OIC-INCAPSULA").grid(row=13, column=0, padx=10, pady=10)

tk.Label(root, text="INCAPSULA Excel File:").grid(row=14, column=0, padx=10, pady=5)
INCAPSULA_excel_entry = tk.Entry(root, width=50)
INCAPSULA_excel_entry.grid(row=14, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(INCAPSULA_excel_entry)).grid(row=14, column=2, padx=10, pady=5)

tk.Label(root, text="INCAPSULA PPTX File:").grid(row=15, column=0, padx=10, pady=5)
INCAPSULA_pptx_entry = tk.Entry(root, width=50)
INCAPSULA_pptx_entry.grid(row=15, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(INCAPSULA_pptx_entry)).grid(row=15, column=2, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_incapsula).grid(row=16, column=0, columnspan=3, pady=20)

#---------------------------------OIC-PALOALTO---------------------------------#
tk.Label(root, text="OIC-PALOALTO").grid(row=17, column=0, padx=10, pady=10)

tk.Label(root, text="PALOALTO Excel File:").grid(row=18, column=0, padx=10, pady=5)
PA_excel_entry = tk.Entry(root, width=50)
PA_excel_entry.grid(row=18, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(PA_excel_entry)).grid(row=18, column=2, padx=10, pady=5)

tk.Label(root, text="PALOALTO PPTX File:").grid(row=19, column=0, padx=10, pady=5)
PA_pptx_entry = tk.Entry(root, width=50)
PA_pptx_entry.grid(row=19, column=1, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(PA_pptx_entry)).grid(row=19, column=2, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_paloalto).grid(row=20, column=0, columnspan=3, pady=20)

#---------------------------------OIC-TRENDMICRO---------------------------------#
tk.Label(root, text="OIC-TRENDMICRO").grid(row=0, column=4, padx=10, pady=10)

tk.Label(root, text="TRENDMICRO Excel File:").grid(row=1, column=4, padx=10, pady=5)
TRENDMICRO_excel_entry = tk.Entry(root, width=50)
TRENDMICRO_excel_entry.grid(row=1, column=5, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(TRENDMICRO_excel_entry)).grid(row=1, column=6, padx=10, pady=5)

tk.Label(root, text="TRENDMICRO PPTX File:").grid(row=2, column=4, padx=10, pady=5)
TRENDMICRO_pptx_entry = tk.Entry(root, width=50)
TRENDMICRO_pptx_entry.grid(row=2, column=5, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(TRENDMICRO_pptx_entry)).grid(row=2, column=6, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_trendmicro).grid(row=3, column=4, columnspan=3, pady=20)

#---------------------------------OIC-VPN---------------------------------#
tk.Label(root, text="OIC-VPN").grid(row=4, column=4, padx=10, pady=10)

tk.Label(root, text="Authentication Success Excel File:").grid(row=5, column=4, padx=10, pady=5)
VPN_Success_excel_entry = tk.Entry(root, width=50)
VPN_Success_excel_entry.grid(row=5, column=5, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(VPN_Success_excel_entry)).grid(row=5, column=6, padx=10, pady=5)

tk.Label(root, text="Authentication Failed File:").grid(row=6, column=4, padx=10, pady=5)
VPN_Failed_excel_entry = tk.Entry(root, width=50)
VPN_Failed_excel_entry.grid(row=6, column=5, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(VPN_Failed_excel_entry)).grid(row=6, column=6, padx=10, pady=5)

tk.Label(root, text="VPN PPTX File:").grid(row=7, column=4, padx=10, pady=5)
VPN_pptx_entry = tk.Entry(root, width=50)
VPN_pptx_entry.grid(row=7, column=5, padx=10, pady=5)
tk.Button(root, text="Browse", command=lambda: open_file_dialog(VPN_pptx_entry)).grid(row=7, column=6, padx=10, pady=5)

tk.Button(root, text="Process Data", command=process_data_vpn).grid(row=8, column=4, columnspan=3, pady=20)

#---------------------------------RUNNING MAIN LOOP---------------------------------#
root.mainloop()



